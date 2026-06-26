"""Office 文档解析器

支持：
- DOCX (python-docx)
- PPTX (python-pptx)
- DOC / PPT → LibreOffice 转换 → 再解析
"""
import subprocess
from pathlib import Path

from app.services.parsers.base import ParsedDocument


def parse_docx(path: str, filename: str) -> ParsedDocument:
    """解析 DOCX → Markdown"""
    from docx import Document as DocxDocument

    doc = DocxDocument(path)
    blocks = []

    # 段落
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            style = paragraph.style.name.lower() if paragraph.style else ""
            if "heading" in style:
                level = style.replace("heading", "").strip()
                try:
                    blocks.append(f"{'#' * int(level)} {text}")
                except ValueError:
                    blocks.append(text)
            else:
                blocks.append(text)

    # 表格
    for table_index, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            blocks.append(f"## Table {table_index}\n\n" + "\n".join(rows))

    title = Path(filename).stem
    content = f"# {title}\n\n" + "\n\n".join(blocks)

    return ParsedDocument(
        filename=filename,
        title=title,
        file_type="docx",
        content=content,
        source_path=path,
        metadata={"source_type": "docx"},
    )


def parse_pptx(path: str, filename: str) -> ParsedDocument:
    """解析 PPTX → Markdown（每页幻灯片为一个分区）"""
    from pptx import Presentation

    presentation = Presentation(path)
    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            # 尝试读取表格
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(cell.text.strip() for cell in row.cells))

        slide_text = "\n\n".join(texts)
        if slide_text:
            # 提取第一个文本作为幻灯片标题
            lines = slide_text.split("\n")
            first_line = lines[0][:80]
            slides.append(f"## Slide {index}: {first_line}\n\n{slide_text}")

    title = Path(filename).stem
    content = f"# {title}\n\n" + "\n\n".join(slides)

    return ParsedDocument(
        filename=filename,
        title=title,
        file_type="pptx",
        content=content,
        source_path=path,
        metadata={
            "source_type": "pptx",
            "slide_count": len(presentation.slides),
        },
    )


def convert_legacy_office(path: str, timeout: int = 120) -> Path:
    """使用 LibreOffice 将旧版 Office 文件转换为新版

    Args:
        path: .doc 或 .ppt 文件路径
        timeout: 转换超时（秒）

    Returns:
        转换后的 .docx 或 .pptx 文件路径
    """
    ext = Path(path).suffix.lower()
    output_ext = ".docx" if ext == ".doc" else ".pptx"
    output_dir = Path(path).parent

    result = subprocess.run(
        [
            "libreoffice",
            "--headless",
            "--convert-to",
            output_ext.lstrip("."),
            "--outdir",
            str(output_dir),
            path,
        ],
        capture_output=True,
        text=True,
        timeout=timeout,
        check=False,
    )

    if result.returncode != 0:
        stderr = result.stderr[-500:] if result.stderr else ""
        raise RuntimeError(f"Office 文件转换失败: {stderr}")

    converted = output_dir / f"{Path(path).stem}{output_ext}"
    if not converted.exists():
        raise RuntimeError(f"转换后文件未生成: {converted}")
    return converted